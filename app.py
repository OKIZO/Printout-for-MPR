import streamlit as st
import json
import io
import os
from pptx import Presentation
from pptx.util import Inches

# レイアウトを横幅いっぱいに使う設定（2カラムに最適化）
st.set_page_config(page_title="PPTX生成システム", layout="wide")

# ==========================================
# 認証・パスワード管理機能
# ==========================================
CONFIG_FILE = "config.json"
DEFAULT_USER_PWD = "team_creative"
ADMIN_PWD = "okino_creative"

# 保存されたパスワードを読み込む関数
def load_user_pwd():
    if os.path.exists(CONFIG_FILE):
        try:
            with open(CONFIG_FILE, "r", encoding="utf-8") as f:
                return json.load(f).get("pwd", DEFAULT_USER_PWD)
        except:
            return DEFAULT_USER_PWD
    return DEFAULT_USER_PWD

# 新しいパスワードを保存する関数
def save_user_pwd(new_pwd):
    with open(CONFIG_FILE, "w", encoding="utf-8") as f:
        json.dump({"pwd": new_pwd}, f)

# セッション状態の初期化
if "logged_in" not in st.session_state:
    st.session_state.logged_in = False

# --- ログイン画面の表示（未ログイン時） ---
if not st.session_state.logged_in:
    # 画面を中央に寄せるためのレイアウト調整
    _, col_center, _ = st.columns([1, 2, 1])
    
    with col_center:
        st.markdown("<h2 style='text-align:center; margin-top:4rem; margin-bottom:2rem;'>🔐 MedConcept ログイン</h2>", unsafe_allow_html=True)
        
        # ログインフォーム
        pwd_input = st.text_input("チーム用パスワードを入力", type="password")
        if st.button("ログイン", type="primary", use_container_width=True):
            if pwd_input == load_user_pwd():
                st.session_state.logged_in = True
                st.rerun() # 画面をリロードしてメインアプリを表示
            else:
                st.error("パスワードが間違っています。")
        
        st.markdown("<br><br>", unsafe_allow_html=True)
        
        # パスワード変更（管理者用）
        with st.expander("⚙️ 管理者設定（パスワードの変更）"):
            st.markdown("<small>※管理者のパスワードが必要です</small>", unsafe_allow_html=True)
            admin_input = st.text_input("管理者パスワード", type="password")
            new_pwd_input = st.text_input("新しいチーム用パスワード")
            
            if st.button("パスワードを更新", use_container_width=True):
                if admin_input == ADMIN_PWD:
                    if new_pwd_input.strip():
                        save_user_pwd(new_pwd_input.strip())
                        st.success(f"チーム用パスワードを「{new_pwd_input.strip()}」に変更しました！")
                    else:
                        st.error("新しいパスワードを入力してください。")
                else:
                    st.error("管理者パスワードが間違っています。")
    
    # ログインしていない場合はここでプログラムを停止し、下のアプリ画面を表示させない
    st.stop()

# ==========================================
# これより下はログイン成功時のみ実行される
# ==========================================

# --- 補助関数：図形やセル内のテキストをフォント維持で置換（分割対策版） ---
def replace_text_in_shape(item, replacements):
    if not hasattr(item, "text_frame") or item.text_frame is None:
        return
    for paragraph in item.text_frame.paragraphs:
        p_text = "".join(run.text for run in paragraph.runs)
        
        replaced_any = False
        for old_text, new_text in replacements.items():
            if old_text in p_text:
                p_text = p_text.replace(old_text, str(new_text))
                replaced_any = True
                
        if replaced_any:
            if len(paragraph.runs) > 0:
                paragraph.runs[0].text = p_text
                for i in range(1, len(paragraph.runs)):
                    paragraph.runs[i].text = ""

# --- 補助関数：不要な図形を完全に削除 ---
def delete_shape(shape):
    try:
        sp_tree = shape.element.getparent()
        sp_tree.remove(shape.element)
    except:
        pass

# --- メイン処理関数 ---
def generate_pptx(json_data, uploaded_images):
    prs = Presentation("template.pptx")

    brand_info = f"カラー：{json_data.get('brandColors', '')}\nブランドイメージ：{'、'.join(json_data.get('brandImages', []))}"
    
    replacements = {
        "{{productName}}": json_data.get("productName", ""),
        "{{itemName}}": json_data.get("itemName", ""),
        "{{spec}}": json_data.get("spec", ""),
        "{{target}}": json_data.get("target", ""),
        "{{scene}}": json_data.get("scene", ""),
        "{{objectiveA}}": json_data.get("objectiveA", ""),
        "{{objectiveB}}": json_data.get("objectiveB", ""),
        "{{before}}": json_data.get("before", ""),
        "{{after}}": json_data.get("after", ""),
        "{{concept}}": json_data.get("concept", ""),
        "{{brandInfo}}": brand_info,
        "{{designExterior}}": "、".join(json_data.get("designExterior", [])),
        "{{functional}}": "、".join(json_data.get("functional", [])),
        "{{toneManner}}": "\n".join(json_data.get("toneManner", [])),
    }

    cb = json_data.get("changeTypesBefore", [])
    ca = json_data.get("changeTypesAfter", [])
    
    for i in range(4):
        replacements[f"{{{{cb{i+1}}}}}"] = cb[i] if i < len(cb) else ""
        replacements[f"{{{{ca{i+1}}}}}"] = ca[i] if i < len(ca) else ""

    for slide in prs.slides:
        def process_shapes(shapes):
            for shape in shapes:
                if shape.shape_type == 6:
                    process_shapes(shape.shapes)
                elif hasattr(shape, "text_frame") and shape.text_frame is not None:
                    replace_text_in_shape(shape, replacements)
                elif shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            replace_text_in_shape(cell, replacements)
        process_shapes(slide.shapes)

    slide_indices = {"A案": 5, "B案": 6, "C案": 7, "D案": 8, "E案": 9}
    margin_x, margin_y = Inches(0.5), Inches(1.5)
    cell_w, cell_h = Inches(3.0), Inches(2.0)
    cols = 3

    for plan_name, images in uploaded_images.items():
        if plan_name in slide_indices and len(prs.slides) > slide_indices[plan_name]:
            slide = prs.slides[slide_indices[plan_name]]
            
            for idx, img_file in enumerate(images[:6]):
                row = idx // cols
                col = idx % cols
                x = margin_x + (col * cell_w)
                y = margin_y + (row * cell_h)
                
                img_stream = io.BytesIO(img_file.read())
                try:
                    slide.shapes.add_picture(img_stream, x, y, width=cell_w - Inches(0.2))
                except Exception as e:
                    st.warning(f"{plan_name}の画像挿入に失敗しました: {e}")

    ppt_stream = io.BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    return ppt_stream

# --- UI構築（左右2カラムレイアウト） ---

# 余白とファイルリストの広がりを抑えるCSS（文字被り修正版）
st.markdown("""
    <style>
        .block-container { padding-top: 1rem; padding-bottom: 1rem; }
        h1 { font-size: 1.6rem !important; margin-bottom: 1rem !important; }
        h2 { font-size: 1.2rem !important; margin-bottom: 0.2rem !important;}
        .stMarkdown p { font-size: 0.85rem; margin-bottom: 0.2rem !important;}
        
        /* ファイルアップローダー周りの隙間を削る */
        [data-testid="stFileUploader"] { margin-bottom: 0rem; }
        
        /* ドロップゾーン（点線の枠）を薄くする */
        [data-testid="stFileUploadDropzone"] {
            padding: 0.5rem !important;
            min-height: 1.5rem !important;
        }
        [data-testid="stFileUploadDropzone"] * {
            font-size: 0.8rem !important;
        }
        [data-testid="stFileUploadDropzone"] svg {
            display: none; /* ドロップゾーンの雲のアイコンを消す */
        }
        
        /* ▼ アップロードされたファイルのリストが下に伸びないようにする ▼ */
        /* リスト全体を小さなスクロール枠に閉じ込める */
        [data-testid="stFileUploader"] > section {
            max-height: 90px !important;  /* 高さを固定（約2ファイル分） */
            overflow-y: auto !important;  /* はみ出たらスクロール */
        }
        
        /* ファイルサイズ（1.2MBなど）の表記を消してスッキリさせる */
        [data-testid="stUploadedFile"] small {
            display: none !important;
        }
    </style>
""", unsafe_allow_html=True)

col1, col2 = st.columns(2, gap="large")

# ===== 左カラム：画像アップロード =====
with col1:
    st.header("🖼️ 画像アップロード")
    st.markdown("各案の画像を枠内にドラッグ＆ドロップしてください。")

    uploaded_images = {}
    plans = ["A案", "B案", "C案", "D案", "E案"]

    for plan in plans:
        # アコーディオンなしで直接アップローダーを表示
        uploaded_images[plan] = st.file_uploader(
            f"📁 {plan}", 
            accept_multiple_files=True, 
            type=["png", "jpg", "jpeg"], 
            key=plan
        )

# ===== 右カラム：JSON入力＆パワポ生成 =====
with col2:
    st.header("📝 企画書生成")
    st.markdown("左側のアプリからコピーしたJSONデータを貼り付けます。")

    # テキストエリアの高さを、左のメニュー群と合うように設定
    json_text = st.text_area("JSONデータを貼り付け", height=280, label_visibility="collapsed", placeholder="ここにJSONデータを貼り付けてください")

    if st.button("📊 企画書パワーポイントを作成", type="primary", use_container_width=True):
        if not json_text.strip():
            st.error("エラー: JSONデータを入力してください。")
        else:
            try:
                json_data = json.loads(json_text)
                with st.spinner("PowerPointを生成中..."):
                    ppt_stream = generate_pptx(json_data, uploaded_images)
                    
                st.success("🎉 PowerPointの生成が完了しました！")
                st.download_button(
                    label="📥 企画書(.pptx) をダウンロード",
                    data=ppt_stream,
                    file_name=f"proposal_{json_data.get('itemName', 'untitled')}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True
                )
                
            except json.JSONDecodeError:
                st.error("エラー: JSONのフォーマットが正しくありません。")
            except Exception as e:
                st.error(f"予期せぬエラーが発生しました: {e}")
